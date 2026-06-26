"""
Document parsers — T14 + T15.
Each parser receives raw bytes and returns plain text + basic metadata.
"""
from __future__ import annotations

import csv
import io
from dataclasses import dataclass, field
from typing import Protocol


@dataclass
class ParsedDocument:
    text: str
    metadata: dict = field(default_factory=dict)
    pages: int = 1


class Parser(Protocol):
    def parse(self, data: bytes, filename: str) -> ParsedDocument: ...


# ── PDF ───────────────────────────────────────────────────────────────────────

class PDFParser:
    def parse(self, data: bytes, filename: str) -> ParsedDocument:
        import fitz  # pymupdf

        doc = fitz.open(stream=data, filetype="pdf")
        page_count = doc.page_count
        pages_text: list[str] = []
        for page in doc:
            pages_text.append(page.get_text("text"))
        text = "\n\n".join(pages_text)
        metadata = {
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
            "subject": doc.metadata.get("subject", ""),
            "pages": page_count,
            "filename": filename,
        }
        doc.close()
        return ParsedDocument(text=text, metadata=metadata, pages=page_count)


# ── DOCX ──────────────────────────────────────────────────────────────────────

class DOCXParser:
    def parse(self, data: bytes, filename: str) -> ParsedDocument:
        from docx import Document

        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Include text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)
        text = "\n\n".join(paragraphs)
        props = doc.core_properties
        metadata = {
            "title": props.title or "",
            "author": props.author or "",
            "filename": filename,
        }
        return ParsedDocument(text=text, metadata=metadata)


# ── PPTX ──────────────────────────────────────────────────────────────────────

class PPTXParser:
    def parse(self, data: bytes, filename: str) -> ParsedDocument:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        slides_text: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_lines: list[str] = [f"[Slide {i}]"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            slide_lines.append(line)
            slides_text.append("\n".join(slide_lines))
        text = "\n\n".join(slides_text)
        metadata = {"filename": filename, "slides": len(prs.slides)}
        return ParsedDocument(text=text, metadata=metadata, pages=len(prs.slides))


# ── XLSX ──────────────────────────────────────────────────────────────────────

class XLSXParser:
    def parse(self, data: bytes, filename: str) -> ParsedDocument:
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sheet_count = len(wb.sheetnames)
        sheets_text: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = [f"[Sheet: {sheet_name}]"]
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(v) for v in row if v is not None)
                if row_text.strip():
                    rows.append(row_text)
            sheets_text.append("\n".join(rows))
        wb.close()
        return ParsedDocument(
            text="\n\n".join(sheets_text),
            metadata={"filename": filename, "sheets": sheet_count},
        )


# ── CSV ───────────────────────────────────────────────────────────────────────

class CSVParser:
    def parse(self, data: bytes, filename: str) -> ParsedDocument:
        text_data = data.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(text_data))
        rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]
        return ParsedDocument(
            text="\n".join(rows),
            metadata={"filename": filename, "rows": len(rows)},
        )


# ── TXT / Markdown ────────────────────────────────────────────────────────────

class TextParser:
    def parse(self, data: bytes, filename: str) -> ParsedDocument:
        text = data.decode("utf-8", errors="replace")
        return ParsedDocument(text=text, metadata={"filename": filename})


# ── HTML ──────────────────────────────────────────────────────────────────────

class HTMLParser:
    def parse(self, data: bytes, filename: str) -> ParsedDocument:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(data, "html.parser")
        # Remove script/style noise
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        title = soup.title.string.strip() if soup.title and soup.title.string else ""
        text = soup.get_text(separator="\n", strip=True)
        return ParsedDocument(text=text, metadata={"filename": filename, "title": title})


# ── Web Scraper (T15) ─────────────────────────────────────────────────────────

class WebScraper:
    def __init__(self, timeout: int = 30) -> None:
        self._timeout = timeout

    def scrape(self, url: str) -> ParsedDocument:
        import httpx
        from bs4 import BeautifulSoup

        resp = httpx.get(url, timeout=self._timeout, follow_redirects=True, headers={
            "User-Agent": "PyRAG-Core/1.0 (+https://github.com/pyrag-core)"
        })
        resp.raise_for_status()
        soup = BeautifulSoup(resp.content, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "aside"]):
            tag.decompose()
        title = soup.title.string.strip() if soup.title and soup.title.string else url
        text = soup.get_text(separator="\n", strip=True)
        return ParsedDocument(
            text=text,
            metadata={"source_url": url, "title": title, "content_type": resp.headers.get("content-type", "")},
        )


# ── Registry ──────────────────────────────────────────────────────────────────

_PARSERS: dict[str, Parser] = {
    "pdf":  PDFParser(),
    "docx": DOCXParser(),
    "pptx": PPTXParser(),
    "xlsx": XLSXParser(),
    "xls":  XLSXParser(),
    "csv":  CSVParser(),
    "txt":  TextParser(),
    "md":   TextParser(),
    "markdown": TextParser(),
    "html": HTMLParser(),
    "htm":  HTMLParser(),
}

SUPPORTED_EXTENSIONS = set(_PARSERS.keys())


def get_parser(extension: str) -> Parser:
    from app.core.exceptions import UnsupportedFileTypeError
    ext = extension.lower().lstrip(".")
    parser = _PARSERS.get(ext)
    if parser is None:
        raise UnsupportedFileTypeError(f"Unsupported file type: .{ext}")
    return parser


def parse_document(data: bytes, filename: str) -> ParsedDocument:
    ext = filename.rsplit(".", 1)[-1] if "." in filename else ""
    parser = get_parser(ext)
    return parser.parse(data, filename)
